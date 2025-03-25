""" Ref for slide types:  
0 ->  title and subtitle 
1 ->  title and content 
2 ->  section header 
3 ->  two content 
4 ->  Comparison 
5 ->  Title only  
6 ->  Blank 
7 ->  Content with caption 
8 ->  Pic with caption 
"""
from pptx import Presentation  
import os


# Creating presentation object 
root = Presentation() 
  
# Loop through all available slide layouts
for i, layout in enumerate(root.slide_layouts):
    # Creating slide object for the current layout
    slide = root.slides.add_slide(layout)

    # Adding title (if available)
    try:
        slide.shapes.title.text = f"Slide Layout {i}"
    except AttributeError:
        pass

    # Iterate through all placeholders in the slide
    try:
        for j, placeholder in enumerate(slide.placeholders):
            placeholder.text = f"Placeholder {j} in layout {i}"
    except AttributeError:
        pass


# Saving file 
if not os.path.exists("Presentations"):
    os.makedirs("Presentations")
root.save("Presentations/Output_all_layouts.pptx")


for slide in root.slides:
    print(f"slide {slide.slide_id-256}:")
    for shape in slide.placeholders:
        if shape.is_placeholder:
            phf = shape.placeholder_format
            print('%d, %s' % (phf.idx, phf.type))




print("done") 