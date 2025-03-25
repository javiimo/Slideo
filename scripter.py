

from pptx import Presentation
import base64
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shapes, text_data, image_data):
    """
    Iterates through all shapes in a slide and extracts text and image data. It handles group shapes

    Args:
        shapes: A list of shapes or a shapes collection from python-pptx.
        text_data: A list to store the extracted text.
        image_data: A list to store the extracted image data.

    Returns:
        A dictionary containing the extracted text and image data.
    """

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP: #In case there is a group shape, use recursion
            iter_shapes(shape.shapes, text_data, image_data)
        else:
            if shape.has_text_frame: #Extract text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_data.append(run.text)
            elif shape.has_table: #Extract tables
                table_data = "Table:\n"
                for row in shape.table.rows:
                    row_data = "|"
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                row_data += run.text + "|"
                    table_data += row_data + "\n"
                text_data.append(table_data)
            if hasattr(shape, 'image'): #Extract images and set them to base64 in jpeg
                image = shape.image
                image_bytes = image.blob
                if image.ext.lower() != ".jpg":
                    from PIL import Image
                    from io import BytesIO
                    image_pil = Image.open(BytesIO(image_bytes))
                    buffered = BytesIO()
                    image_pil.save(buffered, format="JPEG")
                    image_bytes = buffered.getvalue()

                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                image_data.append(image_base64)
            

    return {"text": text_data, "images": image_data}


def extract_text_and_images(slide):
    """
    Extracts all text and image data from a slide.

    Args:
        slide: The slide object from python-pptx.

    Returns:
        A dictionary containing the extracted text and image data.
    """

    text_data = []
    image_data = []

    dict = iter_shapes(slide.shapes, text_data, image_data)
    
    return dict



def call_groq_vision(prompt, client, image=None, model="llama-3.2-11b-vision-preview"):


    sysprompt = "You are an expert narrator script writer. You will receive information extracted from slides, including text and one image (if there is one in the slide). Your task is to create a compelling and informative script for a narrator that explains the slide's content without directly reading the text. Instead, focus on providing context, insights, and connections between the different elements on the slide. If images are present, refer to them generally in your script (e.g., 'as you can see in the accompanying chart...' or 'the image illustrates...'), if there are none, do not say anything about them. Your script should sound natural, engaging, and suitable for a presentation. DO NOT HALLUCINATE INFORMATION, STICK TO THE INFO IN THE SLIDE."

    messages = [
        {"role": "system", "content": sysprompt},
    ]
    if image:
        messages = [
            {
                "role": "user", 
                "content": [
                    {"type": "text", "text": sysprompt +"\n"+ prompt}, 
                    {
                        "type": "image_url", 
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{image}",
                        },
                    },
                ],
            }
        ]
    else:
        messages.append({"role": "user", "content": prompt})

    completion = client.chat.completions.create(
        model=model,
        messages=messages
    )
    return completion.choices[0].message.content


def scripter(presentation_path, comments, client):
    """
    Extracts information from each slide of a presentation and uses the Groq API to generate a narrator script.

    Args:
        presentation_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of strings, where each string is the narrator script for a slide.
    """

    prs = Presentation(presentation_path)
    scripts = []

    for i, slide in enumerate(prs.slides):
        extracted_data = extract_text_and_images(slide)

        # Construct the prompt for the Groq API
        prompt = f"Slide {i + 1}:\n"
        if extracted_data["text"]:
            prompt += "Text contained in the slide:\n" + "\n".join(extracted_data["text"]) + "\n"
        if comments:
            prompt += "Comments and suggestions about the script:\n" + comments + "\n" #Adding the comments to the prompt
        prompt += """Provide the narrator script directly without any introductory or concluding remarks (e.g., 'Here's the script...') or other additions like stage directions (e.g., '(a brief pause)').  Do not include anything else, just the script.
        
        Adjust the time of narration to the type of slide you are presenting: 
        - For slides that only show a section title, keep it VERY brief and seamless, simply introducing the section. 
        - For the index slide, provide a brief overview of the presentation's structure and main points. Focus on the big picture without going into details for each section. Keep it concise.
        - For content slides, provide more detailed explanations and insights, paraphrasing the information contained in the slide to enhance clarity. 
        Ensure EVERYTHING you say is grounded in the slide's content.
        """

        # Encode the first image if any
        image = None
        if extracted_data["images"]:
            image = extracted_data["images"][0]

        # Call the Groq API to generate the script
        script = call_groq_vision(prompt, client, image)
        scripts.append(script)

    return scripts



if __name__ == "__main__":
    #Testing scripter
    presentation_path = "out/output.pptx"
    comments = "" 
    scripts = scripter(presentation_path, comments)
    with open("out/script.txt", "w") as f:
        for i, script in enumerate(scripts):
            f.write(f"Slide {i+1}:\n{script}\n\n")
    print("Script saved to out/script.txt")




    #Testing extraction
    # presentation_path = "out/out.pptx"
    # prs = Presentation(presentation_path)
    # with open("out/extraction.txt", "w", encoding="utf-8") as f:
    #     for i, slide in enumerate(prs.slides):
    #         extracted_data = extract_text_and_images(slide)
    #         print(f"Slide {i}:", file=f)
    #         if extracted_data["text"]:
    #             print("Text:", file=f)
    #             for text in extracted_data["text"]:
    #                 print(text, file=f)
    #         if extracted_data["images"]:
    #             print("Images:", file=f)
    #             for image in extracted_data["images"]:
    #                 print("THERE IS AN IMAGE HERE!", file=f)  
    #         print(file=f)
    



