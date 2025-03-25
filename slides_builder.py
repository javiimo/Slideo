from pptx import Presentation 
from pptx.util import Inches
import json
import re


def get_measures(image_paths):
    if len(image_paths) > 4: #If there are more than 4 images give, just take the first 4
        image_paths = image_paths[:4]
        print("There are more than 4 images in a slide. We are taking just the first 4.")
    #Extract the sizes from the image paths
    sizes = []
    for image_path in image_paths:
        # Extracting the size from the image filename
        size_str = image_path.split('_')[-1].split('.')[0]  # Get the last part before the extension
        width, height = map(int, size_str.split('x'))  # Split by 'x' and convert to integers
        sizes.append([width, height])
    
    #Compute the ratios
    ratios = [size[1]/size[0] for size in sizes]
    if len(ratios) == 1:
        w1 = min([5/ratios[0],9])
        h1 = ratios[0] * w1
        measures = [
            [
                Inches(0.5 + (9-w1)/2), #left
                Inches(1.6 + (5-h1)/2), #top
                Inches(w1) #width
            ]
        ]
    elif len(ratios) == 2:
        w1 = min([5/ratios[0],4.5])
        h1 = ratios[0] * w1
        w2 = min([5/ratios[1],4.5])
        h2 = ratios[1] * w2
        measures = [
            [
                Inches(0.5 + (4.5-w1)/2), #left
                Inches(1.6 + (5-h1)/2), #top
                Inches(w1) #width
            ],
            [
                Inches(5 + (4.5-w2)/2), #left
                Inches(1.6 + (5-h2)/2), #top
                Inches(w2) #width
            ]
        ]
    elif len(ratios) in [3,4]:
        w1 = min([2.5/ratios[0],4.5])
        h1 = ratios[0] * w1
        w2 = min([2.5/ratios[1],4.5])
        h2 = ratios[1] * w2
        w3 = min([2.5/ratios[2],4.5])
        h3 = ratios[2] * w3
        measures = [
            [
                Inches(0.5 + (4.5-w1)/2), #left
                Inches(1.6 + (2.5-h1)/2), #top
                Inches(w1) #width
            ],
            [
                Inches(5 + (4.5-w2)/2), #left
                Inches(1.6 + (2.5-h2)/2), #top
                Inches(w2) #width
            ],
            [
                Inches(0.5 + (4.5-w3)/2), #left
                Inches(4.1 + + (2.5-h3)/2), #top
                Inches(w3) #width
            ]
        ]
        if len(ratios) == 4:
            w4 = min([2.5/ratios[3],4.5])
            h4 = ratios[3] * w4
            measures.append(
                    [
                        Inches(5 + (4.5-w4)/2), #left
                        Inches(4.1 + + (2.5-h4)/2), #top
                        Inches(w4) #width
                    ]
                )
    else: 
        print(f"The lenght of the image paths is: {len(image_paths)}")
        return
    return measures



def add_one_image(image, slide): #This slide must be from the layout index 8
    placeholder = slide.placeholders[1]  # Get the first shape from the second slide
    ph_w = placeholder.width + Inches(2.6)
    ph_h = placeholder.height
    ph_l = placeholder.left - Inches(1.3)
    ph_t = placeholder.top
    ph_ratio = ph_h / ph_w
    size_str = image.split('_')[-1].split('.')[0]  # Get the last part before the extension
    width, height = map(int, size_str.split('x'))  # Split by 'x' and convert to integers
    img_ratio = height / width
    if ph_ratio > img_ratio:
        slide.shapes.add_picture(image, ph_l, ph_t + (ph_h - ph_w*img_ratio)/2, ph_w)
    else: 
        slide.shapes.add_picture(image, ph_l + (ph_w - ph_h/img_ratio)/2, ph_t, height = ph_h)
    placeholder.element.getparent().remove(placeholder.element)





def build_presentation(indices, images) -> Presentation:
    prs = Presentation()
    for i, j in zip(indices, images):
        if i == 5: #Add multiple images slide
            slide_layout = prs.slide_layouts[i]
            slide = prs.slides.add_slide(slide_layout)
            if j:
                image_paths = j
                measures = get_measures(image_paths)
                for aux, image in enumerate(image_paths):
                    left, top, width = measures[aux]
                    slide.shapes.add_picture(image, left, top, width)

        elif i == 8: #Add single image slide
            slide_layout = prs.slide_layouts[i]
            slide = prs.slides.add_slide(slide_layout)  
            if j:  # Check if j is not an empty list
                image = j[0]
                add_one_image(image, slide)

        elif i == 11: #create the index slide
            slide_layout = prs.slide_layouts[5]
            ind_slide = prs.slides.add_slide(slide_layout)
            ind_title = ind_slide.placeholders[0]
            ind_title.text = "Index"
            left = Inches(1.0)
            top = Inches(1.5) 
            width = Inches(8.0) # adjust width as needed
            height = Inches(5.5) # adjust height as needed
            ind_slide.shapes.add_textbox(left, top, width, height)
        
        else:
            try: 
                slide_layout = prs.slide_layouts[i]
                prs.slides.add_slide(slide_layout)  
            except Exception as e:
                print(f"An error occurred while adding a slide: {str(e)}")
                print("Skipping to the next slide. There is a wrong index in indices")
            
    return prs


def call_structured_groq(prompt, client):
    sysprompt = """You are an expert slide creator specializing in crafting clear, concise, and professional presentations. Your task is to populate slide content based on a provided description, strictly adhering to the format and purpose of each slide type within a presentation.  Prioritize clarity and conciseness in your content creation, ensuring each slide delivers information effectively and maintains a professional tone.

    Consider the following guidelines when generating content:

    * **Accuracy:** Ensure all information presented is factually correct and relevant to the slide's topic.
    * **Brevity:**  Use concise language, avoiding jargon and unnecessary complexity.  Prioritize bullet points and short sentences for easy readability.
    * **Visual Appeal:** Structure content to be visually appealing on a slide, using appropriate headings, subheadings, and bullet points.  Avoid overcrowding slides with excessive text.
    * **Engagement:** Craft content that is engaging and informative, capturing the audience's attention and effectively conveying key messages.
    * **Professionalism:** Maintain a professional tone throughout the presentation, using appropriate language and avoiding informal or colloquial expressions.
    * **Slide Type Adherence:**  Pay close attention to the specified slide type and tailor the content accordingly.  For example, a title slide should only contain the title and subtitle, while a content slide should have a title and supporting points.
    * **Contextual Relevance:** Ensure the generated content aligns with the overall presentation topic and flows logically from previous slides.

    You will be provided with the following information for each slide:

    * **Slide Type:**  The type of slide (e.g., title, content, comparison, image).
    * **Description:** A detailed description of the information to be included on the slide.

    Your output should be structured to fit seamlessly within the designated slide template, providing all necessary text elements in a clear and organized manner. Ensure your generated content is ready to be directly inserted into the corresponding slide placeholder(s).  Do not include any formatting instructions or markup; simply provide the plain text content for the slide."""

    completion = client.chat.completions.create(
        model="llama3-8b-8192",
        messages=[
            {"role": "system", "content": sysprompt},
            {"role": "user", "content": prompt},
        ],
        response_format={"type": "json_object"},
    )
    return completion



"""
0 : Title
2 : Section Title
1 : Content in bullet points
4 : Content comparison in 2 columns
8 : Only one image
5 : More than one image
11 : Index
"""
def add_content(presentation, client, indices, descriptions):
    for i, slide in enumerate(presentation.slides):
        if indices[i] == 0: #Title. Set: title and subtitle
            prompt = f"""
            Create content for a title slide. It contains a title and subtitle for the whole presentation. Write it as a final version, do not write any comments or mention the duration of the presentation.

            Description of this slide: {descriptions[i]}

            Return the title and subtitle as a JSON object with the following schema:
            ```json
            {{"title": "title text", "subtitle": "subtitle text"}}
            ```
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            title = content["title"]
            subtitle = content["subtitle"]
            slide.placeholders[0].text = title
            slide.placeholders[1].text = subtitle

        elif indices[i] == 2: #Section Title. Set: section_title, section_subtitle
            prompt = f"""
            Create content for a section title slide. It contains a title and a subtitle of the section.

            Description of this slide: {descriptions[i]}

            Return the title and subtitle as a JSON object with the following schema:
            ```json
            {{"title": "title text", "subtitle": "subtitle text"}}
            ```
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            section_title = content["title"]
            section_subtitle = content["subtitle"]
            slide.placeholders[0].text = section_title
            slide.placeholders[1].text = section_subtitle

        elif indices[i] == 1: #Content in bullet points. Set: slide_title, bullets
            prompt = f"""
            Create content for a content slide with bullet points. It contains a title and bullet points.

            Description of this slide: {descriptions[i]}

            Return the title and bullet points as a JSON object with the following schema:
            ```json
            {{"title": "title text", "bullets": ["bullet 1", "bullet 2", ...]}}
            ```
            Where the bullets field is a list of strings and each string represents a bullet point.
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            slide_title = content["title"]
            bullets = content["bullets"]
            slide.placeholders[0].text = slide_title
            slide.placeholders[1].text = "\n".join(bullets)

        elif indices[i] == 4: #Content comparison in 2 columns. Set: slide_title, col1 and col2 titles and bullets
            prompt = f"""
            Create content for a comparison slide with 2 columns. It contains a title and two columns with titles and bullet points. The bullets should be employed to compare the 2 concepts that each column represent.

            Description of this slide: {descriptions[i]}

            Return the title, columns titles and bullet points as a JSON object with the following schema:
            ```json
            {{"title": "title text", "col1_title": "col1_title", "col1_bullets": ["bullet 1", "bullet 2", ...], "col2_title": "col2_title", "col2_bullets": ["bullet 1", "bullet 2", ...]}}
            ```
            Where the col1_bullets and col2_bullets fields are lists of strings and each string represents a bullet point.
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            slide_title = content["title"]
            col1_title = content["col1_title"]
            col1_bullets = content["col1_bullets"]
            col2_title = content["col2_title"]
            col2_bullets = content["col2_bullets"]
            slide.placeholders[0].text = slide_title
            slide.placeholders[1].text = col1_title
            slide.placeholders[2].text = "\n".join(col1_bullets)
            slide.placeholders[3].text = col2_title
            slide.placeholders[4].text = "\n".join(col2_bullets)

        elif indices[i] == 8: #Only one image. Set: slide_title, img_caption, img_subcaption
            prompt = f"""
            Create content for a picture slide. It only contains the image caption and a subcaption for you to fill.

            Description of this slide: {descriptions[i]}

            Return the title, image caption and image subcaption as a JSON object with the following schema:
            ```json
            {{"img_caption": "img_caption text", "img_subcaption": "img_subcaption text"}}
            ```
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            img_caption = content["img_caption"]
            img_subcaption = content["img_subcaption"]
            slide.placeholders[0].text = img_caption
            slide.placeholders[2].text = img_subcaption

        elif indices[i] == 5: #More than one image
            prompt = f"""
            Create content for a picture slide. It contains only a title for you to fill. DO NOT INCLUDE ANY COMMENTS OR SUGGESTIONS IN THE TITLE, KEEP IT VERY CONCISE AND RELEVANT FOR THE SLIDE IMAGES.

            Description of this slide: {descriptions[i]}

            Return the title, image captions and image subcaptions as a JSON object with the following schema:
            ```json
            {{"title": "title text"}}
            ```
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            slide_title = content["title"]
            slide.placeholders[0].text = slide_title + " (include images in this slide)"

        elif indices[i] == 11: #Index. Set: index
            prompt = f"""
            Create content for an index slide. It contains a list of sections with their respective hierarchical numbering. FIX ANY TYPOS, FORMAT ERRORS, OR DUPLICATIONS OF SECTIONS YOU MIGHT FIND IN THE DESCRIPTION OF THE INDEX. Also, do not include in the index any comments that are found in the description but shouldn't be in a final index.

            Unformatted description of the index draft of the presentation: {descriptions[i]}

            Return a JSON object with the following schema.  The "index" key should contain a list of strings. Each string represents a section title prefixed with its hierarchical numbering. Subsections are indicated by increasing the decimal numbering. For instance, "1.1" is a subsection under "1", and "1.1.1" is a sub-subsection under "1.1".

            ```json
            {{
              "index": [
                "1. First section title",
                "1.1 First subsection title",
                "1.2 Second subsection title",
                "2. Second section title",
                "2.1 First subsection title",
                ...
              ]
            }}
            ```

            IMPORTANT: Don't include the words "section" or "subsection" in the titles themselves, just the numbering and the name. Ensure the numbering accurately reflects the hierarchical structure described in the provided presentation index description.

            Example:

            Given the description:
            "Slide 1: Introduction. Slide 2: Methods. Slide 3: Results. Slide 4: First experiment. Slide 5: Second experiment. Slide 6: Discussion and conclusions."

            The correct JSON output would be:
            ```json
            {{
              "index": [
                "1. Introduction",
                "2. Methods",
                "3. Results",
                "3.1 First experiment",
                "3.2 Second experiment",
                "4. Discussion and conclusions"
              ]
            }}
            ```

            Note how section 3 has subsections ("3.1" and "3.2") while others don't.  Carefully follow the provided description to infer this structure.  Invalid JSON or a schema mismatch will cause an error.
            Use Arabic numerals (1, 2, 3...) for numbering the sections and subsections in the index. Do not use Roman numerals (I, II, III...).
            """

            completion = call_structured_groq(prompt, client)
            content = json.loads(completion.choices[0].message.content)
            index = content["index"]
            slide.placeholders[0].text = "Table of Contents" #Setting the title of the slide
            

            #This paragraph sets the correct indentation according to the numbering
            text_frame = slide.shapes[1].text_frame
            for para_str in index:
                p = text_frame.add_paragraph()
                match = re.match(r"(\d+(\.\d+)*)", para_str[:10])
                if match:
                    p.level = len(match.group(1).split(".")) -1
                p.text = para_str
        else:
            print(f"There is a wrong index: {indices[i]}. Skipping to the next slide.")
    return presentation






if __name__ == "__main__":
    indices = [0, 11, 1, 2, 1, 4, 5, 1, 8, 1]
    descriptions = ['Title Slide: The Rise of Smart Beauty Devices with FOREO', 'Here is the slide index:\n\nI. Introduction to Smart Beauty Devices (Slide 1: "Introduction to Smart Beauty Devices")\n\t* Title Slide: The Rise of Smart Beauty Devices with FOREO (Slide 0)\n\nII. Market Trends and Consumer Behavior (Slide 2: "Market Trends and Consumer Behavior")\n\nIII. Features and Benefits of Smart Beauty Devices (Slide 3: "Key Features and Benefits of Smart Beauty Devices")\n\t* Timer: 45 seconds\n\nIV. Comparison of FOREO Devices (Slides 4-5)\n\t* Slide 4: "Comparison Slide: FOREO Luna vs. FOREO Issa"\n\t* Image: Slide 5: "Picture Slide: FOREO Devices Comparison Picture"\n\t* Timer: 45 seconds\n\nV. Benefits of Using Smart Beauty Devices (Slide 6: "Benefits of Using Smart Beauty Devices")\n\t* Timer: 45 seconds\n\nVI. FOREO Device Illustrations (Slide 7: "Picture Slide: Face Cleansing Device from FOREO")\n\t* Image: "face_cleansing_device"\n\t* Timer: 30 seconds\n\nVII. Conclusion and Future Outlook (Slide 8: "Conclusion and Future Outlook")\n\t* Timer: 45 seconds', 'Section Title: Introduction to Smart Beauty Devices\nSection of context about the beauty industry and trend of smart beauty devices', 'Section Title: Market Trends and Consumer Behavior\nSection title only, no content', 'Content: Key Features and Benefits of Smart Beauty Devices\nBullet points explaining key features and benefits of smart beauty devices\nTimer: 45 seconds', 'Comparison Slide: FOREO Luna vs. FOREO Issa\nTwo columns comparing FOREO Luna and FOREO Issa devices, with key features and benefits\nTimer: 45 seconds', 'Picture Slide: FOREO Devices Comparison Picture\nBoth FOREO Luna and FOREO Issa devices side-by-side for comparison, image: foreo_devices_comparison_picture\nTimer: 30 seconds', 'Content: Benefits of Using Smart Beauty Devices\nBullet points explaining benefits of using smart beauty devices\nTimer: 45 seconds', 'Picture Slide: Face Cleansing Device from FOREO\nSingle image of face cleansing device from FOREO, image: face_cleansing_device\nTimer: 30 seconds', 'Content: Conclusion and Future Outlook\nBullet points concluding the presentation and looking forward to the future of smart beauty devices\nTimer: 45 seconds']


    images = [
            [],
            [],
            [],
            [],
            [],
            [],
            [r"img\FOREO smart beauty devices comparison\FOREO_LUNA_device_..._1280x720.jpg"],
            [],
            [r"img\smart beauty gadgets FOREO widgets images\FOREO_Launches_LUNA_fofo_-_The_Worlds_..._1440x640.jpg"],
            []
        ]
    prs = build_presentation(indices, images)
    #add_content(prs, indices, descriptions)

    prs.save("out/try.pptx")
